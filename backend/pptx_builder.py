"""Minimal PPTX builder for generating presentation outputs from structured content."""
from __future__ import annotations

from pathlib import Path
from typing import Any

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
except ImportError:
    Presentation = None


def pptx_available() -> bool:
    return Presentation is not None


def create_title_slide(prs: Any, title: str, subtitle: str = "") -> None:
    layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = title
    if subtitle and slide.placeholders[1]:
        slide.placeholders[1].text = subtitle


def create_content_slide(prs: Any, title: str, bullets: list[str]) -> None:
    layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = title
    body = slide.placeholders[1]
    tf = body.text_frame
    tf.clear()
    for i, bullet in enumerate(bullets):
        if i == 0:
            tf.paragraphs[0].text = bullet
        else:
            p = tf.add_paragraph()
            p.text = bullet
            p.level = 0


def create_image_slide(prs: Any, title: str, image_path: str, caption: str = "") -> None:
    layout = prs.slide_layouts[5]  # blank
    slide = prs.slides.add_slide(layout)
    # Add title textbox
    from pptx.util import Inches, Pt
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    tf = txBox.text_frame
    tf.text = title
    tf.paragraphs[0].font.size = Pt(24)
    tf.paragraphs[0].font.bold = True
    # Add image
    img_path = Path(image_path)
    if img_path.exists():
        slide.shapes.add_picture(str(img_path), Inches(1), Inches(1.3), Inches(8), Inches(5))
    if caption:
        cap_box = slide.shapes.add_textbox(Inches(1), Inches(6.5), Inches(8), Inches(0.5))
        cap_box.text_frame.text = caption
        cap_box.text_frame.paragraphs[0].font.size = Pt(10)


def create_table_slide(prs: Any, title: str, headers: list[str], rows: list[list[str]]) -> None:
    layout = prs.slide_layouts[5]  # blank
    slide = prs.slides.add_slide(layout)
    from pptx.util import Inches, Pt
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    tf = txBox.text_frame
    tf.text = title
    tf.paragraphs[0].font.size = Pt(24)
    tf.paragraphs[0].font.bold = True

    num_cols = len(headers)
    num_rows = len(rows) + 1
    col_width = Inches(9.0 / num_cols)
    table_shape = slide.shapes.add_table(num_rows, num_cols, Inches(0.5), Inches(1.3), Inches(9), Inches(5))
    table = table_shape.table

    for i, header in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = header
        cell.text_frame.paragraphs[0].font.bold = True
        cell.text_frame.paragraphs[0].font.size = Pt(10)

    for row_idx, row in enumerate(rows):
        for col_idx, val in enumerate(row):
            cell = table.cell(row_idx + 1, col_idx)
            cell.text = str(val)
            cell.text_frame.paragraphs[0].font.size = Pt(9)


def build_report_pptx(
    output_path: str | Path,
    title: str,
    subtitle: str,
    sections: list[dict[str, Any]],
) -> Path:
    """Build a PPTX from structured sections.
    
    Each section dict can have:
      - type: "content" | "image" | "table"
      - title: slide title
      - bullets: list[str] (for content)
      - image_path: str (for image)
      - caption: str (for image)
      - headers: list[str] (for table)
      - rows: list[list[str]] (for table)
    """
    if not pptx_available():
        raise RuntimeError("python-pptx is not installed")

    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    create_title_slide(prs, title, subtitle)

    for section in sections:
        slide_type = section.get("type", "content")
        slide_title = section.get("title", "")
        if slide_type == "content":
            create_content_slide(prs, slide_title, section.get("bullets", []))
        elif slide_type == "image":
            create_image_slide(prs, slide_title, section.get("image_path", ""), section.get("caption", ""))
        elif slide_type == "table":
            create_table_slide(prs, slide_title, section.get("headers", []), section.get("rows", []))

    out = Path(output_path)
    out.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out))
    return out
